import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import io


# =========================================================
# FUNGSI BIKIN REKOMENDASI PER PIC
# =========================================================
def build_rekomendasi(pic_table: pd.DataFrame):
    """
    Menghasilkan dict:
    { PIC: [list kalimat rekomendasi] }
    berdasarkan tabel kinerja PIC.
    """
    rekom_map = {}

    for _, row in pic_table.iterrows():
        pic = row["PIC"]
        avg_age = row["AVG_AGE"]
        sla = row["SLA_%"]
        open_cnt = row["OPEN"]
        closed_cnt = row["CLOSED"]

        rekom = []

        if avg_age > 60:
            rekom.append(
                "- Ticket lama > **60 hari**. Perlu daily follow-up & koordinasi lintas divisi."
            )

        if sla < 70:
            rekom.append(
                "- SLA penyelesaian < **70%**. Perlu penambahan ritme closing ticket dan monitoring ketat."
            )

        if open_cnt > closed_cnt:
            rekom.append(
                "- Jumlah ticket OPEN lebih banyak dari CLOSED → potensi backlog, perlu prioritas penyelesaian."
            )

        if len(rekom) == 0:
            rekom.append("✔ Performance sangat baik & stabil. Pertahankan pola kerja saat ini.")

        rekom_map[pic] = rekom

    return rekom_map


# =========================================================
# FUNGSI GENERATE PPT (PAKAI TEMPLATE MORATEL)
# =========================================================
def generate_ppt(
    df_full: pd.DataFrame,
    status_summary: pd.DataFrame,
    age_summary: pd.DataFrame,
    pic_monthly: pd.DataFrame,
    pic_perf_table: pd.DataFrame,
    rekom_map: dict,
    top_aging: pd.DataFrame,
    template_path: str = "template PPT Moratel.pptx",
):
    """
    Menghasilkan PPT bytes berdasarkan data full.
    Tidak memakai gambar chart, hanya teks & bullet supaya aman di Streamlit Cloud.
    """
    try:
        prs = Presentation(template_path)
    except Exception as e:
        raise FileNotFoundError(
            f"Template PPT '{template_path}' tidak ditemukan di repo: {e}"
        )

    # ------------- SLIDE 1: COVER -------------
    title_slide_layout = prs.slide_layouts[0]  # biasanya Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    total_ticket = len(df_full)
    now_str = datetime.now().strftime("%d %B %Y")

    title.text = "Laporan Ticket CODEX"
    subtitle.text = f"Generated otomatis dari Dashboard Streamlit\nTotal Ticket: {total_ticket}\nTanggal: {now_str}"

    # ------------- SLIDE 2: RINGKASAN STATUS -------------
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Ringkasan Status Ticket"

    body = slide.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True

    total_critical = df_full["STATUS"].str.contains("critical", case=False).sum()
    total_warning = df_full["STATUS"].str.contains("warning", case=False).sum()

    p = body.paragraphs[0]
    p.text = f"Total Ticket: {total_ticket}"
    p.level = 0

    p = body.add_paragraph()
    p.text = f"Total Critical: {total_critical}"
    p.level = 0

    p = body.add_paragraph()
    p.text = f"Total Warning: {total_warning}"
    p.level = 0

    p = body.add_paragraph()
    p.text = "Rincian per kategori STATUS:"
    p.level = 0

    # detail status
    for _, row in status_summary.iterrows():
        p = body.add_paragraph()
        p.text = f"- {row['STATUS']}: {row['JUMLAH']} ticket (Avg Age: {row['AVG_AGE']} hari, Max: {row['MAX_AGE']} hari)"
        p.level = 1

    # ------------- SLIDE 3: RINGKASAN UMUR TICKET -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Ringkasan Umur Ticket (AGE_DAYS)"

    body = slide.placeholders[1].text_frame
    body.clear()

    overall_avg_age = round(df_full["AGE_DAYS"].mean(), 1)
    overall_max_age = int(df_full["AGE_DAYS"].max())

    p = body.paragraphs[0]
    p.text = f"Average Age semua ticket: {overall_avg_age} hari"
    p.level = 0

    p = body.add_paragraph()
    p.text = f"Ticket tertua: {overall_max_age} hari"
    p.level = 0

    p = body.add_paragraph()
    p.text = "Distribusi kategori umur ticket:"
    p.level = 0

    # bikin bucket umur
    bins = [0, 30, 60, 90, 180, 365, 9999]
    labels = ["0–30", "31–60", "61–90", "91–180", "181–365", ">365"]
    df_full["AGE_BUCKET"] = pd.cut(df_full["AGE_DAYS"], bins=bins, labels=labels, right=True)
    bucket = df_full["AGE_BUCKET"].value_counts().sort_index()

    for bucket_label, val in bucket.items():
        p = body.add_paragraph()
        p.text = f"- {bucket_label} hari: {val} ticket"
        p.level = 1

    # ------------- SLIDE 4: PRODUKTIVITAS PIC PER BULAN -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Produktivitas PIC per Bulan"

    body = slide.placeholders[1].text_frame
    body.clear()

    # batasi agar tidak kepanjangan di slide
    p = body.paragraphs[0]
    p.text = "Ringkasan jumlah ticket per PIC per bulan (max 20 baris):"
    p.level = 0

    pic_monthly_sorted = (
        pic_monthly.sort_values(["PIC", "MONTH"])
        .head(20)
        .reset_index(drop=True)
    )

    for _, row in pic_monthly_sorted.iterrows():
        p = body.add_paragraph()
        p.text = f"- {row['PIC']} — {row['MONTH']}: {row['JUMLAH']} ticket"
        p.level = 1

    # ------------- SLIDE 5: PERFORMANCE PIC (SLA) -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Performance PIC (Total / Open / Close / SLA)"

    body = slide.placeholders[1].text_frame
    body.clear()

    p = body.paragraphs[0]
    p.text = "Ringkasan kinerja per PIC (max 15 baris):"
    p.level = 0

    perf_sorted = pic_perf_table.sort_values("TOTAL", ascending=False).head(15)

    for _, row in perf_sorted.iterrows():
        p = body.add_paragraph()
        p.text = (
            f"- {row['PIC']}: TOTAL={row['TOTAL']} | "
            f"OPEN={row['OPEN']} | CLOSED={row['CLOSED']} | "
            f"SLA={row['SLA_%']}% | AVG_AGE={row['AVG_AGE']} hari"
        )
        p.level = 1

    # ------------- SLIDE 6: REKOMENDASI PER PIC -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Rekomendasi Perbaikan Kinerja PIC"

    body = slide.placeholders[1].text_frame
    body.clear()

    first = True
    for pic, rekom_list in rekom_map.items():
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = f"{pic}:"
        p.level = 0

        for r in rekom_list:
            p = body.add_paragraph()
            # hapus markdown ** di versi PPT
            clean_r = r.replace("**", "")
            p.text = clean_r
            p.level = 1

    # ------------- SLIDE 7: TOP 10 TICKET PALING TUA -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Top 10 Ticket Paling Tua"

    body = slide.placeholders[1].text_frame
    body.clear()

    p = body.paragraphs[0]
    p.text = "Daftar 10 ticket dengan AGE_DAYS tertinggi:"
    p.level = 0

    for _, row in top_aging.iterrows():
        p = body.add_paragraph()
        p.text = (
            f"- {row['NO-TICKET']} | PIC={row['PIC']} | STATUS={row['STATUS_TICKET']} | "
            f"AGE={row['AGE_DAYS']} hari"
        )
        p.level = 1

    # ------------- SLIDE 8: CATATAN / CLOSING -------------
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Catatan & Tindak Lanjut"

    body = slide.placeholders[1].text_frame
    body.clear()

    p = body.paragraphs[0]
    p.text = "Catatan umum:"
    p.level = 0

    p = body.add_paragraph()
    p.text = "- Ticket dengan umur tinggi perlu menjadi prioritas mingguan dalam rapat ENG/NOC/BB/Access."
    p.level = 1

    p = body.add_paragraph()
    p.text = "- SLA PIC bisa dijadikan dasar penentuan beban kerja dan kebutuhan tambahan resource."
    p.level = 1

    p = body.add_paragraph()
    p.text = "- Dashboard ini bisa dijalankan berkala (harian/mingguan) sebagai bahan laporan manajemen."
    p.level = 1

    # Simpan ke bytes
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io


# =========================================================
# STREAMLIT APP
# =========================================================
st.set_page_config(page_title="Dashboard Ticket CODEX", layout="wide")

st.title("📊 Dashboard Analisa Ticket CODEX — Versi Lengkap")

uploaded = st.file_uploader("📁 Upload File CODEX (.xlsx)", type=["xlsx"])
if uploaded is None:
    st.info("Silakan upload file CODEx (.xlsx) terlebih dahulu.")
    st.stop()

# -------------------------
# LOAD & PREPARE DATA FULL
# -------------------------
df = pd.read_excel(uploaded)
df.columns = [c.strip() for c in df.columns]

required_cols = [
    "NO-TICKET",
    "HOSTNAME",
    "INTERFACE",
    "STATUS",
    "ASSIGN DIVISION",
    "CREATE TICKET",
]

for col in required_cols:
    if col not in df.columns:
        st.error(f"Kolom **{col}** tidak ditemukan. Periksa file Excel.")
        st.stop()

df["CREATE TICKET"] = pd.to_datetime(df["CREATE TICKET"], errors="coerce")
df["AGE_DAYS"] = (datetime.now() - df["CREATE TICKET"]).dt.days

df["SUB_DIVISI"] = df["ASSIGN DIVISION"].str.split("-").str[-1].str.strip()
df["PIC"] = df["ASSIGN DIVISION"].str.split("-").str[0].str.strip()

if "STATUS.1" in df.columns:
    df["STATUS_TICKET"] = df["STATUS.1"]
else:
    df["STATUS_TICKET"] = df["STATUS"]

# -------------------------
# SIDEBAR FILTER
# -------------------------
st.sidebar.header("⚙️ Filter Data")

subdivisi_opt = ["ALL"] + sorted(df["SUB_DIVISI"].dropna().unique().tolist())
status_opt = ["ALL"] + sorted(df["STATUS_TICKET"].dropna().unique().tolist())
pic_opt = ["ALL"] + sorted(df["PIC"].dropna().unique().tolist())

bulan_opt = ["ALL"] + sorted(df["CREATE TICKET"].dt.month_name().dropna().unique().tolist())
tahun_opt = ["ALL"] + sorted(df["CREATE TICKET"].dt.year.dropna().unique().tolist())

f_sub = st.sidebar.selectbox("Sub Divisi", subdivisi_opt)
f_status = st.sidebar.selectbox("Status Ticket", status_opt)
f_pic = st.sidebar.selectbox("PIC", pic_opt)
f_bulan = st.sidebar.selectbox("Bulan", bulan_opt)
f_tahun = st.sidebar.selectbox("Tahun", tahun_opt)

filtered = df.copy()

if f_sub != "ALL":
    filtered = filtered[filtered["SUB_DIVISI"] == f_sub]
if f_status != "ALL":
    filtered = filtered[filtered["STATUS_TICKET"] == f_status]
if f_pic != "ALL":
    filtered = filtered[filtered["PIC"] == f_pic]
if f_bulan != "ALL":
    filtered = filtered[filtered["CREATE TICKET"].dt.month_name() == f_bulan]
if f_tahun != "ALL":
    filtered = filtered[filtered["CREATE TICKET"].dt.year == f_tahun]

st.success(f"Total data setelah filter: {len(filtered)}")

# =========================================================
# SUMMARY NUMBER CARD
# =========================================================
st.subheader("📌 Ringkasan Data (berdasarkan filter aktif)")
c1, c2, c3, c4 = st.columns(4)

c1.metric("Total Ticket", len(filtered))
c2.metric(
    "Critical",
    filtered["STATUS"].str.contains("Critical", case=False, na=False).sum(),
)
c3.metric(
    "Warning",
    filtered["STATUS"].str.contains("Warning", case=False, na=False).sum(),
)
c4.metric("Average Age (Days)", round(filtered["AGE_DAYS"].mean(), 1))

# =========================================================
# GRAFIK STATUS
# =========================================================
st.subheader("📊 Distribusi Ticket Berdasarkan Kategori STATUS")

count_status = filtered["STATUS"].value_counts().reset_index()
count_status.columns = ["STATUS", "JUMLAH"]

fig_status = px.bar(
    count_status,
    x="STATUS",
    y="JUMLAH",
    color="STATUS",
    color_discrete_sequence=px.colors.qualitative.Set1,
)
st.plotly_chart(fig_status, use_container_width=True)

# =========================================================
# TABEL RINCIAN STATUS
# =========================================================
st.subheader("📄 Tabel Rincian Status Ticket")

summary = (
    filtered.groupby("STATUS")
    .agg(
        JUMLAH=("STATUS", "count"),
        AVG_AGE=("AGE_DAYS", "mean"),
        MAX_AGE=("AGE_DAYS", "max"),
    )
    .reset_index()
)

summary["AVG_AGE"] = summary["AVG_AGE"].round(1)
st.dataframe(summary, use_container_width=True)

# =========================================================
# HISTOGRAM AGE DAYS
# =========================================================
st.subheader("⏳ Distribusi Umur Ticket (AGE_DAYS)")

fig_age = px.histogram(
    filtered,
    x="AGE_DAYS",
    nbins=30,
    color="STATUS",
    color_discrete_sequence=px.colors.qualitative.Set2,
)
st.plotly_chart(fig_age, use_container_width=True)

# =========================================================
# PRODUKTIVITAS PIC PER BULAN (FILTERED)
# =========================================================
st.subheader("🧑‍💻 Performance PIC per Bulan (berdasarkan filter)")

filtered["MONTH"] = filtered["CREATE TICKET"].dt.to_period("M").astype(str)
pic_perf_filtered = (
    filtered.groupby(["PIC", "MONTH"])
    .size()
    .reset_index(name="JUMLAH")
)

fig_pic = px.bar(
    pic_perf_filtered,
    x="PIC",
    y="JUMLAH",
    color="MONTH",
    barmode="group",
)
st.plotly_chart(fig_pic, use_container_width=True)

# =========================================================
# TABEL KINERJA PIC LENGKAP (FILTERED)
# =========================================================
st.subheader("📋 Tabel Performance PIC (berdasarkan filter)")

pic_table = (
    filtered.groupby("PIC")
    .agg(
        TOTAL=("PIC", "count"),
        OPEN=("STATUS_TICKET", lambda x: (x == "Open").sum()),
        CLOSED=("STATUS_TICKET", lambda x: (x == "Close").sum()),
        AVG_AGE=("AGE_DAYS", "mean"),
        MAX_AGE=("AGE_DAYS", "max"),
    )
    .reset_index()
)

pic_table["SLA_%"] = round((pic_table["CLOSED"] / pic_table["TOTAL"]) * 100, 1)
pic_table["AVG_AGE"] = pic_table["AVG_AGE"].round(1)

st.dataframe(pic_table, use_container_width=True)

# =========================================================
# REKOMENDASI OTOMATIS (FILTERED)
# =========================================================
st.subheader("💡 Rekomendasi Perbaikan Kinerja PIC")

rekom_map_filtered = build_rekomendasi(pic_table)

for pic, rekom_list in rekom_map_filtered.items():
    st.markdown(f"### 🔧 {pic}")
    for r in rekom_list:
        st.write(r)
    st.write(
        "**Alasan bisnis:** Penyelesaian ticket cepat mengurangi risiko alarm berulang, "
        "mengurangi downtime, dan mempercepat troubleshooting NOC/BB/Access.\n---"
    )

# =========================================================
# DATA LENGKAP (FILTERED)
# =========================================================
st.subheader("📑 Data Lengkap Ticket CODEX (berdasarkan filter)")
st.dataframe(filtered, use_container_width=True)

# =========================================================
# BAGIAN DOWNLOAD PPT – DATA FULL (TANPA FILTER)
# =========================================================
st.subheader("📥 Download Laporan PPT (Template Moratel, Data FULL)")

# --- siapkan data FULL untuk PPT ---
df_full = df.dropna(subset=["CREATE TICKET"]).copy()

# summary status full
status_summary_full = (
    df_full.groupby("STATUS")
    .agg(
        JUMLAH=("STATUS", "count"),
        AVG_AGE=("AGE_DAYS", "mean"),
        MAX_AGE=("AGE_DAYS", "max"),
    )
    .reset_index()
)
status_summary_full["AVG_AGE"] = status_summary_full["AVG_AGE"].round(1)

# summary age (dipakai untuk slide, tapi sebagian sudah dihitung di generate_ppt)
age_summary_full = status_summary_full[["STATUS", "AVG_AGE", "MAX_AGE"]].copy()

# produktivitas PIC per bulan full
df_full["MONTH"] = df_full["CREATE TICKET"].dt.to_period("M").astype(str)
pic_monthly_full = (
    df_full.groupby(["PIC", "MONTH"])
    .size()
    .reset_index(name="JUMLAH")
)

# kinerja PIC full
pic_perf_full = (
    df_full.groupby("PIC")
    .agg(
        TOTAL=("PIC", "count"),
        OPEN=("STATUS_TICKET", lambda x: (x == "Open").sum()),
        CLOSED=("STATUS_TICKET", lambda x: (x == "Close").sum()),
        AVG_AGE=("AGE_DAYS", "mean"),
        MAX_AGE=("AGE_DAYS", "max"),
    )
    .reset_index()
)
pic_perf_full["SLA_%"] = round((pic_perf_full["CLOSED"] / pic_perf_full["TOTAL"]) * 100, 1)
pic_perf_full["AVG_AGE"] = pic_perf_full["AVG_AGE"].round(1)

# rekomendasi full (supaya konsisten antara dashboard & PPT)
rekom_map_full = build_rekomendasi(pic_perf_full)

# top 10 aging
top_aging_full = (
    df_full.sort_values("AGE_DAYS", ascending=False)
    .head(10)[["NO-TICKET", "PIC", "STATUS_TICKET", "AGE_DAYS"]]
)

# generate PPT bytes
try:
    ppt_bytes = generate_ppt(
        df_full=df_full,
        status_summary=status_summary_full,
        age_summary=age_summary_full,
        pic_monthly=pic_monthly_full,
        pic_perf_table=pic_perf_full,
        rekom_map=rekom_map_full,
        top_aging=top_aging_full,
        template_path="template PPT Moratel.pptx",
    )

    st.download_button(
        label="⬇ Download PPT Laporan Lengkap (Data FULL)",
        data=ppt_bytes,
        file_name="Laporan_Ticket_CODEX_Moratel.pptx",
        mime=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )
except FileNotFoundError as e:
    st.error(
        "Template PPT tidak ditemukan di repo. "
        "Pastikan file bernama **'template PPT Moratel.pptx'** ada di root repository."
    )
    st.text(str(e))
except Exception as e:
    st.error(f"Terjadi error saat generate PPT: {e}")
